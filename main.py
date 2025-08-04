# main.py

# Imports de bibliotecas padrão
import os
import csv
import json
import tempfile
import logging
import platform
import re

# Imports de bibliotecas externas
import pymysql
import requests
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor

# Importa as variáveis de configuração do nosso arquivo config.py
from config import (
    DB_HOST, DB_USER, DB_PASSWORD, DB_NAME,
    TRAINING_ASSETS_BASE_PATH, MOVIDESK_API_TOKEN,
    MOVIDESK_VERSION_FIELD_ID, MOVIDESK_OTHER_FIELD_ID, MOVIDESK_OTHER_FIELD_RULE_ID,
    MOVIDESK_OWNER_ID, MOVIDESK_OWNER_TEAM_NAME, MOVIDESK_ACTION_CREATOR_ID,
    ACTION_HTML_SIGNATURE, validar_configuracoes
)


# --- CONFIGURAÇÕES E SETUP INICIAL ---

LOGS_DIR = "logs"
RELATORIOS_DIR = "relatorios"
PPTX_DIR = "powerpoint"
PPTX_TEMPLATE_PATH = "Layout-Base.pptx"

class Cores:
    """Classe para armazenar códigos de cores ANSI para o terminal."""
    RESET = '\033[0m'
    VERDE = '\033[92m'
    AMARELO = '\033[93m'
    VERMELHO = '\033[91m'
    AZUL = '\033[94m'
    CIANO = '\033[96m'

def criar_diretorios():
    """Cria os diretórios de output se não existirem."""
    print("Verificando e criando diretórios de saída...")
    os.makedirs(LOGS_DIR, exist_ok=True)
    os.makedirs(RELATORIOS_DIR, exist_ok=True)
    os.makedirs(PPTX_DIR, exist_ok=True)
    print(f"Diretórios '{LOGS_DIR}/', '{RELATORIOS_DIR}/', '{PPTX_DIR}/' prontos.")

def limpar_tela():
    """Limpa a tela do console."""
    os.system('cls' if platform.system() == 'Windows' else 'clear')

def exibir_cabecalho(titulo):
    """Exibe um cabeçalho formatado para cada seção."""
    limpar_tela()
    print(Cores.AZUL + "="*60)
    print(f"   {titulo.upper()}")
    print("="*60 + Cores.RESET)
    print()

class SuccessFilter(logging.Filter):
    def filter(self, record):
        return record.levelno == logging.INFO and 'distribuído com sucesso' in record.getMessage()

class ErrorFilter(logging.Filter):
    def filter(self, record):
        return record.levelno >= logging.ERROR

# Configuração do Logger
logger = logging.getLogger('dist')
logger.setLevel(logging.DEBUG)
fh_all = logging.FileHandler(os.path.join(LOGS_DIR, 'distribuicao_tickets.log'), 'w', 'utf-8')
fh_all.setLevel(logging.DEBUG)
fh_all.setFormatter(logging.Formatter('%(asctime)s | %(levelname)s | %(message)s'))
logger.addHandler(fh_all)
fh_succ = logging.FileHandler(os.path.join(LOGS_DIR, 'tickets_success.txt'), 'w', 'utf-8')
fh_succ.setLevel(logging.INFO)
fh_succ.addFilter(SuccessFilter())
fh_succ.setFormatter(logging.Formatter('%(message)s'))
logger.addHandler(fh_succ)
fh_err = logging.FileHandler(os.path.join(LOGS_DIR, 'tickets_error.txt'), 'w', 'utf-8')
fh_err.setLevel(logging.ERROR)
fh_err.addFilter(ErrorFilter())
fh_err.setFormatter(logging.Formatter('%(message)s'))
logger.addHandler(fh_err)

# --- FUNÇÕES DE MANIPULAÇÃO DE IMAGEM E PPTX ---
def add_order_to_image(image_path, order, total, margin=20):
    """Adiciona texto de ordem no canto inferior direito da imagem."""
    with Image.open(image_path) as img:
        if img.mode != 'RGB': img = img.convert('RGB')
        draw = ImageDraw.Draw(img)
        font_size = max(20, int(max(img.width, img.height) * 0.01))
        try: font = ImageFont.truetype("arial.ttf", font_size)
        except IOError: font = ImageFont.load_default()
        text = f"{order}/{total}".strip()
        text_bbox = draw.textbbox((0, 0), text, font=font)
        text_width = text_bbox[2] - text_bbox[0]
        text_height = text_bbox[3] - text_bbox[1]
        x = max(margin, img.width - text_width - margin)
        y = max(margin, img.height - text_height - margin)
        padding = 5
        draw.rectangle([(x - padding, y), (x + text_width + padding, y + text_height + padding)], fill=(0, 0, 0))
        draw.text((x, y), text, font=font, fill=(255, 255, 255))
        temp_file = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        img.save(temp_file.name)
        return temp_file.name

def add_slide_with_title(prs, layout, title, version, tickets_data):
    """Adiciona um slide com título e versão."""
    slide = prs.slides.add_slide(layout)
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == 1:
            slide.shapes._spTree.remove(shape._element)
    title_placeholder = slide.shapes.title
    title_placeholder.text = f"{version}"
    title_placeholder.text_frame.paragraphs[0].font.size = Pt(20)
    title_placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 109, 105)
    slide_width = prs.slide_width
    text_width = slide_width * 0.65
    margin_left = slide_width * 0.21
    margin_right = slide_width - (margin_left + text_width)
    if margin_right < slide_width * 0.1:
        margin_right = slide_width * 0.1
        text_width = slide_width - (margin_left + margin_right)
    textbox = slide.shapes.add_textbox(margin_left, Inches(-0.275), text_width, Inches(1))
    text_frame = textbox.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = False
    p.font.name = "Roboto"
    text_frame.word_wrap = True
    add_footer_texts(slide, title, tickets_data, prs)
    return slide

def add_footer_texts(slide, titulo, tickets_data, prs):
    """Adiciona os textos de rodapé ao slide."""
    parts = titulo.split(' - ', 1)
    suite_id = parts[0].strip()
    ticket_info = tickets_data.get(suite_id, {"movidesk": "N/A", "ordem": "N/A"})
    movidesk_text = f"Movidesk: {ticket_info['movidesk']}"
    left_textbox = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.6), Inches(4), Inches(0.5))
    left_text_frame = left_textbox.text_frame
    left_text_frame.clear()
    p_left = left_text_frame.add_paragraph()
    p_left.text = movidesk_text
    p_left.font.size = Pt(12)
    p_left.font.name = "Roboto"
    p_left.alignment = PP_ALIGN.LEFT
    left_text_frame.word_wrap = True

def add_images_with_animation(prs, layout, title, version, img_paths, tickets_data):
    """Adiciona múltiplos prints e aplica a sequência de animação."""
    slide = add_slide_with_title(prs, layout, title, version, tickets_data)
    if not img_paths: return
    slide_width, slide_height = prs.slide_width, prs.slide_height
    total_images = len(img_paths)
    image_shapes = []
    for idx, img_path in enumerate(img_paths):
        temp_image_path = None
        try:
            temp_image_path = add_order_to_image(img_path, idx + 1, total_images)
            pic = slide.shapes.add_picture(temp_image_path, 0, 0)
            max_width, max_height = slide_width * 0.8, slide_height * 0.7
            pic.width, pic.height = _calculate_new_dimensions(pic.width, pic.height, max_width, max_height)
            pic.left, pic.top = (slide_width - pic.width) // 2, (slide_height - pic.height) // 2
            if idx > 0: pic.element.spPr.append(parse_xml('<a:noFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'))
            image_shapes.append(pic)
        finally:
            if temp_image_path and os.path.exists(temp_image_path): os.unlink(temp_image_path)
    if total_images <= 1: return
    main_sequence = _get_or_create_main_sequence(slide)
    id_counter = IdCounter(10)
    for i in range(total_images - 1):
        exit_node = _create_animation_node(image_shapes[i].shape_id, 'exit', 'clickEffect', id_counter)
        appear_node = _create_animation_node(image_shapes[i+1].shape_id, 'appear', 'afterEffect', id_counter)
        main_sequence.append(_create_click_group(exit_node, appear_node, id_counter))

def _calculate_new_dimensions(img_width, img_height, max_width, max_height):
    ratio = min(max_width / img_width, max_height / img_height) if img_width > 0 and img_height > 0 else 0
    return int(img_width * ratio), int(img_height * ratio)

def _get_or_create_main_sequence(slide):
    timing = slide.element.find('.//p:timing', slide.element.nsmap)
    if timing is None:
        timing = parse_xml('<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:tnLst><p:par><p:cTn id="1" dur="indefinite" nodeType="tmRoot"><p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst/></p:cTn><p:prevCondLst><p:cond evt="onPrev"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst><p:nextCondLst><p:cond evt="onNext"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst></p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>')
        slide.element.append(timing)
    return timing.find('.//p:cTn[@nodeType="mainSeq"]/p:childTnLst', slide.element.nsmap)

def _create_click_group(exit_anim_node, appear_anim_node, id_counter):
    click_group_node = parse_xml(f'<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="{id_counter.get()}" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{id_counter.increment()}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst/></p:cTn></p:par></p:childTnLst></p:cTn></p:par>')
    id_counter.increment()
    anim_list_parent = click_group_node.find('.//p:cTn/p:childTnLst/p:par/p:cTn/p:childTnLst', click_group_node.nsmap)
    anim_list_parent.append(exit_anim_node)
    anim_list_parent.append(appear_anim_node)
    return click_group_node

def _create_animation_node(shape_id, effect, node_type, id_counter):
    visibility = "visible" if effect == 'appear' else "hidden"
    preset = "entr" if effect == 'appear' else "exit"
    anim_node = parse_xml(f'<p:par xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"><p:cTn id="{id_counter.get()}" fill="hold" nodeType="{node_type}" presetClass="{preset}" presetID="1" presetSubtype="0"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id="{id_counter.get()+1}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val="{visibility}"/></p:to></p:set></p:childTnLst></p:cTn></p:par>')
    id_counter.increment(2)
    return anim_node

class IdCounter:
    def __init__(self, start_id): self._id = start_id
    def get(self): return self._id
    def increment(self, value=1): self._id += value; return self._id

# --- FUNÇÕES DE LÓGICA DE NEGÓCIO E API ---
def read_tickets_csv(csv_path):
    """Lê o arquivo CSV e retorna um dicionário com os dados."""
    tickets_data = {}
    if not os.path.exists(csv_path): return tickets_data
    with open(csv_path, mode='r', encoding='utf-8') as file:
        reader = csv.DictReader(file, delimiter=';')
        for row in reader:
            if suite_id := row.get('suite', '').strip():
                tickets_data[suite_id] = {k: v.strip() for k, v in row.items()}
    return tickets_data

def process_directory(base_dir, version, prs, layout_new, layout_old, version_novo, mode, tickets_data, numero_final):
    """Processa diretórios para montar os slides."""
    counts = {"new": 0, "old": 0}
    for section in ["new", "old"]:
        section_dir = os.path.join(base_dir, section)
        if not os.path.exists(section_dir): continue
        layout = layout_new if section == "new" else layout_old
        sorted_folders = ordenar_pastas([p for p in os.listdir(section_dir) if os.path.isdir(os.path.join(section_dir, p))], numero_final)
        for folder in sorted_folders:
            folder_path = os.path.join(section_dir, folder)
            folder_num = extrair_numero(folder, numero_final)
            if (mode == "F" and folder.lower() != "novo") or \
                (mode == "L" and (folder.lower() == "novo" or folder_num > numero_final)) or \
                (mode == "P" and folder.lower() != "novo" and folder_num <= numero_final):
                continue
            current_version = version_novo if folder.lower() == "novo" else (version if folder.lower() == "final" else f"{'.'.join(version.split('.')[:2])}.{int(folder.split(' ')[-1]):03d}")
            task_folders = sorted([(tickets_data.get(f.split(' - ')[0].strip(), {}).get("ordem", 999999), f) for f in os.listdir(folder_path) if os.path.isdir(os.path.join(folder_path, f))])
            for _, task_folder in task_folders:
                task_path = os.path.join(folder_path, task_folder)
                prints = sorted([os.path.join(task_path, f) for f in os.listdir(task_path) if f.lower().endswith(('.png', '.jpg', '.jpeg'))])
                if prints:
                    add_images_with_animation(prs, layout, task_folder, current_version, prints, tickets_data)
                    counts[section] += 1
    print(f"\n{Cores.AZUL}{'='*60}\n      RESUMO DO PROCESSAMENTO DA APRESENTAÇÃO\n{'='*60}{Cores.RESET}")
    print(f"Total de novas implementações (new) processadas: {counts['new']}")
    print(f"Total de old (old) processados: {counts['old']}")
    print(f"{Cores.AZUL}{'='*60}{Cores.RESET}")

def extrair_numero(pasta, numero_final):
    """Extrai o número de uma pasta para ordenação."""
    if pasta.lower() == "final": return numero_final
    if pasta.lower() == "novo": return float('inf')
    match = re.search(r'\d+', pasta)
    return int(match.group()) if match else float('inf') - 1

def ordenar_pastas(pastas, numero_final):
    """Ordena as pastas com base no número extraído."""
    return sorted(pastas, key=lambda p: extrair_numero(p, numero_final))

def distribute_tickets(tickets_path, version):
    """Lê o CSV e distribui os tickets no Movidesk."""
    if not os.path.exists(tickets_path):
        print(f"{Cores.VERMELHO}ERRO: Arquivo {tickets_path} não encontrado.{Cores.RESET}")
        return
    with open(tickets_path, mode='r', encoding='utf-8') as file:
        reader = csv.DictReader(file, delimiter=';')
        for row in reader:
            if movidesk_num := row.get('movidesk'):
                post_movidesk(movidesk_num, version, row.get('observacao', ''))

def get_ticket_details(ticket_id, token):
    """Busca detalhes de um ticket na API do Movidesk."""
    url = f"https://api.movidesk.com/public/v1/tickets"
    params = {'token': token, '$filter': f"id eq {ticket_id}", '$select': 'id,customFieldValues', '$expand': 'customFieldValues($expand=items)'}
    try:
        response = requests.get(url, params=params)
        response.raise_for_status()
        data = response.json()
        items = data.get('items', data)
        return items[0] if isinstance(items, list) and items else (data if isinstance(data, dict) and 'id' in data else None)
    except requests.RequestException as e:
        logger.error(f"Erro na API ao buscar ticket {ticket_id}. Erro: {e}")
        return None

def post_movidesk(ticket_id, version, observacao=''):
    """Atualiza um ticket no Movidesk."""
    ticket_details = get_ticket_details(ticket_id, MOVIDESK_API_TOKEN)
    custom_fields = [{"items": [], "customFieldId": MOVIDESK_VERSION_FIELD_ID, "customFieldRuleId": 620, "line": 1, "value": version}]
    if ticket_details:
        for field in ticket_details.get('customFieldValues', []):
            if field.get('customFieldId') == MOVIDESK_OTHER_FIELD_ID and field.get('customFieldRuleId') == MOVIDESK_OTHER_FIELD_RULE_ID:
                custom_fields.append({"customFieldId": MOVIDESK_OTHER_FIELD_ID, "customFieldRuleId": MOVIDESK_OTHER_FIELD_RULE_ID, "value": field.get('value'), "line": field.get('line', 1)})
                break
    descricao = f"Versão <strong>{version}</strong> distribuída!🚀<br>Qualquer dúvida, estamos à disposição! 💬"
    if observacao: descricao += f"<br><br><p><strong>Observação:</strong></p> {observacao}"
    descricao += f"<div class=\"signature\">{ACTION_HTML_SIGNATURE}</div>"
    payload = {
        "status": "Nova Version", "justification": "",
        "owner": {"id": MOVIDESK_OWNER_ID, "personType": 1, "profileType": 2},
        "ownerTeam": MOVIDESK_OWNER_TEAM_NAME,
        "customFieldValues": custom_fields,
        "actions": [{"type": 1, "origin": 9, "description": descricao, "status": "Nova Version", "createdBy": {"id": MOVIDESK_ACTION_CREATOR_ID, "personType": 2, "profileType": 2}}]
    }
    try:
        response = requests.patch(f"https://api.movidesk.com/public/v1/tickets?token={MOVIDESK_API_TOKEN}&id={ticket_id}", json=payload)
        response.raise_for_status()
        logger.info(f'{ticket_id} – distribuído com sucesso')
    except requests.RequestException:
        logger.error(f'{ticket_id} – erro na distribuição', exc_info=True)

def find_task_folder_by_id(base_dir, suite_id):
    """Encontra a pasta de uma tarefa pelo ID."""
    for section in ["new", "old"]:
        section_dir = os.path.join(base_dir, section)
        if not os.path.exists(section_dir): continue
        for folder in os.listdir(section_dir):
            if os.path.isdir(os.path.join(section_dir, folder)):
                for task_folder in os.listdir(os.path.join(section_dir, folder)):
                    if task_folder.strip().startswith(str(suite_id).strip()):
                        return {"path": os.path.join(section_dir, folder, task_folder)}
    return None

def has_image_files(directory_path):
    """Verifica se um diretório contém arquivos de imagem."""
    if not os.path.isdir(directory_path): return False
    return any(f.lower().endswith(('.png', '.jpg', '.jpeg')) for f in os.listdir(directory_path))

def verificar_projeto_no_banco(project_id, base_dir):
    """Verifica pendências de um projeto e gera um relatório."""
    print(f"\n--- Iniciando verificação do Projeto ID: {project_id} ---")
    ID_DO_SEU_PAPEL_QA = 2 # Exemplo de ID de papel de negócio
    problemas = []
    try:
        with pymysql.connect(host=DB_HOST, user=DB_USER, password=DB_PASSWORD, database=DB_NAME, charset='utf8mb4') as connection:
            with connection.cursor() as cursor:
                cursor.execute(f"SELECT umUsuario FROM papelDoUser WHERE identificador = {project_id} AND identificador2 = {ID_DO_SEU_PAPEL_QA}")
                QAs_do_projeto = {row[0] for row in cursor.fetchall()}
                cursor.execute(f"SELECT tf.task, tf.tk, tf.tt, tf.or FROM utft AS tf JOIN u_tk AS tk ON tf.tk = tk.ntk WHERE tk.identificador = {project_id}")
                for task, nid_ticket, titulo, ordem in cursor.fetchall():
                    if not nid_ticket: continue
                    QA_name = None
                    if QAs_do_projeto:
                        in_clause = f"IN {tuple(QAs_do_projeto)}" if len(QAs_do_projeto) > 1 else f"= {list(QAs_do_projeto)[0]}"
                        cursor.execute(f"SELECT u.uname FROM ucom AS c JOIN usua AS u ON c.idusu = u.umUsuario WHERE c.task = {task} AND c.idusu {in_clause} ORDER BY c.tinc ASC LIMIT 1")
                        if result := cursor.fetchone(): QA_name = result[0]
                    problema = None
                    if not (folder_info := find_task_folder_by_id(base_dir, nid_ticket)): problema = "PASTA NÃO ENCONTRADA"
                    elif not has_image_files(folder_info["path"]): problema = "PASTA CRIADA, MAS SEM PRINTS"
                    if problema: problemas.append({"tk": nid_ticket, "ordem": ordem or 999999, "titulo": titulo, "QA": QA_name, "problema": problema})
    except pymysql.MySQLError as e:
        print(f"{Cores.VERMELHO}ERRO DE BANCO DE DADOS: {e}{Cores.RESET}")
        return
    output_file = os.path.join(RELATORIOS_DIR, f"relatorio_verificacao_projeto_{project_id}.txt")
    if not problemas:
        print(f"\n{Cores.VERDE}VERIFICAÇÃO CONCLUÍDA: Nenhum problema encontrado.{Cores.RESET}")
        if os.path.exists(output_file): os.remove(output_file)
        return
    print(f"\nVERIFICAÇÃO CONCLUÍDA: {len(problemas)} problemas encontrados. Gerando relatório...")
    problemas_agrupados = {}
    for p in problemas: (problemas_agrupados.setdefault(p["QA"] or "TAREFAS SEM QA", [])).append(p)
    with open(output_file, 'w', encoding='utf-8') as f:
        for QA in sorted(problemas_agrupados.keys(), key=lambda k: (k == "TAREFAS SEM QA", k)):
            f.write(f"{QA.upper()}\n\n")
            for p in sorted(problemas_agrupados[QA], key=lambda x: x["ordem"]):
                f.write(f'{p["tk"]} / {p["ordem"] if p["ordem"] != 999999 else "S/O"} - {p["titulo"]} -> {p["problema"]}\n')
            f.write("\n")
    print(f"{Cores.VERDE}Relatório gerado com sucesso em '{output_file}'.{Cores.RESET}")

def generate_csv_from_project(project_id, version_string, output_file):
    """Gera um CSV com todos os tickets de um projeto."""
    print(f"\n--- Iniciando geração de CSV para o Projeto ID: {project_id} ---")
    resultados = []
    try:
        with pymysql.connect(host=DB_HOST, user=DB_USER, password=DB_PASSWORD, database=DB_NAME, charset='utf8mb4') as connection:
            with connection.cursor() as cursor:
                cursor.execute(f"SELECT tf.tk, tf.tt, tf.or FROM utft AS tf JOIN u_tk AS tk ON tf.tk = tk.ntk WHERE tk.identificador = {project_id} ORDER BY tf.or ASC")
                for suite_id, titulo, ordem in cursor.fetchall():
                    if not suite_id: continue
                    cursor.execute(f"SELECT resumo FROM unc WHERE ntk = '{suite_id}' ORDER BY dt_altera ASC LIMIT 1")
                    movidesk_num = "Não encontrado"
                    if result := cursor.fetchone():
                        try: movidesk_num = result[0].split("Movidesk:")[1].split()[0]
                        except (IndexError, AttributeError): pass
                    resultados.append({"suite": suite_id, "titulo": titulo, "movidesk": movidesk_num, "ordem": ordem or 999999, "observacao": ""})
    except pymysql.MySQLError as e:
        print(f"{Cores.VERMELHO}ERRO DE BANCO DE DADOS: {e}{Cores.RESET}")
        return
    with open(output_file, mode='w', newline='', encoding='utf-8') as file:
        writer = csv.DictWriter(file, fieldnames=["suite", "titulo", "movidesk", "ordem", "observacao"], delimiter=';')
        writer.writeheader()
        writer.writerows(sorted(resultados, key=lambda x: x["ordem"]))
    print(f"\n{Cores.AZUL}{'='*50}\nRESUMO DA GERAÇÃO - Projeto {project_id}\nTotal de tickets processados: {len(resultados)}\nArquivo gerado: {output_file}\n{'='*50}{Cores.RESET}")

# --- INTERFACE INTERATIVA (CLI) ---

def get_input(prompt, example="", is_numeric=False):
    """Solicita uma entrada do usuário e a valida."""
    novo_prompt = f"{Cores.AMARELO}>> {prompt} {Cores.CIANO}({example}): {Cores.RESET}"
    while not (user_input := input(novo_prompt).strip()) or (is_numeric and not user_input.isdigit()):
        print(f"{Cores.VERMELHO}Entrada inválida. Tente novamente.{Cores.RESET}")
    return user_input

def get_ppt_mode():
    """Solicita e valida o modo de geração do PowerPoint."""
    print(f"{Cores.AMARELO}>> Escolha o modo de geração do PowerPoint:{Cores.RESET}")
    options = {"F": "Apenas a pasta 'novo'", "L": "Apenas as pastas 'Letra'", "A": "Ambos (Letras e novo)", "P": "Pós-final (Após final até novo)"}
    for key, value in options.items(): print(f"  [{Cores.CIANO}{key}{Cores.RESET}] - {value}")
    while (choice := input(f"{Cores.CIANO}Opção: {Cores.RESET}").strip().upper()) not in options:
        print(f"{Cores.VERMELHO}Opção inválida. Escolha uma das opções acima.{Cores.RESET}")
    return choice

def main():
    """Função principal que executa o menu interativo."""
    criar_diretorios()
    logger.debug('Logger inicializado.')
    while True:
        exibir_cabecalho("FERRAMENTA DE AUTOMAÇÃO DE TREINAMENTOS E TICKETS")
        print("\nEscolha uma opção:")
        print(f"  {Cores.VERDE}1.{Cores.RESET} Verificar Pendências de Projeto")
        print(f"  {Cores.VERDE}2.{Cores.RESET} Gerar Arquivo CSV de Tickets por Projeto")
        print(f"  {Cores.VERDE}3.{Cores.RESET} Gerar Apresentação PowerPoint (.pptx)")
        print(f"  {Cores.VERDE}4.{Cores.RESET} Distribuir Tickets no Movidesk")
        print(f"\n  {Cores.VERMELHO}5. Sair{Cores.RESET}")
        choice = input(f"\n{Cores.AMARELO}>> Digite o número da opção desejada: {Cores.RESET}").strip()

        if choice == '1':
            exibir_cabecalho("1. VERIFICAR PENDÊNCIAS")
            project_id = get_input("Digite o ID do Projeto", "Ex: 115432", is_numeric=True)
            version_novo = get_input("Digite a Versão novo", "Ex: v123")
            base_dir = os.path.join(TRAINING_ASSETS_BASE_PATH, version_novo)
            verificar_projeto_no_banco(project_id, base_dir)

        elif choice == '2':
            exibir_cabecalho("2. GERAR ARQUIVO CSV")
            project_id = get_input("Digite o ID do Projeto", "Ex: 115432", is_numeric=True)
            version_novo = get_input("Digite a Versão novo", "Ex: v123")
            output_file = os.path.join(RELATORIOS_DIR, f"TicketsTreinamento_Projeto_{project_id}.csv")
            generate_csv_from_project(project_id, version_novo, output_file)

        elif choice == '3':
            exibir_cabecalho("3. GERAR APRESENTAÇÃO POWERPOINT")
            tickets_path = "TicketsTreinamento.csv"
            print(f"{Cores.AMARELO}Atenção:{Cores.RESET} Esta função requer um arquivo chamado '{Cores.CIANO}{tickets_path}{Cores.RESET}'.")
            if not os.path.exists(tickets_path):
                print(f"\n{Cores.VERMELHO}ERRO: Arquivo '{tickets_path}' não encontrado!{Cores.RESET}")
            else:
                version_novo = get_input("Digite a Versão novo", "Ex: v123")
                version_final_str = get_input("Digite a Versão final", "Ex: BIGL")
                try: numero_final = int(version_final_str.split('.')[-1])
                except (IndexError, ValueError): numero_final = 10
                mode = get_ppt_mode()
                base_dir = os.path.join(TRAINING_ASSETS_BASE_PATH, version_novo)
                try:
                    prs = Presentation(PPTX_TEMPLATE_PATH)
                    tickets_data = read_tickets_csv(tickets_path)
                    process_directory(base_dir, version_final_str, prs, prs.slide_layouts[3], prs.slide_layouts[2], version_novo, mode, tickets_data, numero_final)
                    output_file_name = f"treinamento_{mode}_{version_novo}_{version_final_str}.pptx"
                    output_path = os.path.join(PPTX_DIR, output_file_name)
                    prs.save(output_path)
                    print(f"\n{Cores.VERDE}Apresentação salva como '{output_path}' com sucesso!{Cores.RESET}")
                except Exception as e:
                    print(f"\n{Cores.VERMELHO}Ocorreu um erro ao gerar a apresentação: {e}{Cores.RESET}", exc_info=True)

        elif choice == '4':
            exibir_cabecalho("4. DISTRIBUIR TICKETS NO MOVIDESK")
            tickets_path = "TicketsTreinamento_Distribuicao.csv"
            print(f"{Cores.AMARELO}Atenção:{Cores.RESET} Esta função requer um arquivo chamado '{Cores.CIANO}{tickets_path}{Cores.RESET}'.")
            if not os.path.exists(tickets_path):
                print(f"\n{Cores.VERMELHO}ERRO: Arquivo '{tickets_path}' não encontrado!{Cores.RESET}")
            else:
                version_to_distribute = get_input("Digite a versão que será informada na distribuição", "Ex: v123")
                distribute_tickets(tickets_path, version_to_distribute)
                print(f"\n{Cores.VERDE}Distribuição concluída. Verifique os arquivos de log.{Cores.RESET}")

        elif choice == '5':
            print(f"\n{Cores.CIANO}Encerrando o programa. Até logo!{Cores.RESET}")
            break
        else:
            print(f"\n{Cores.VERMELHO}Opção inválida. Por favor, escolha um número de 1 a 5.{Cores.RESET}")

        input(f"\n{Cores.AMARELO}Pressione Enter para voltar ao menu...{Cores.RESET}")

if __name__ == "__main__":
    validar_configuracoes()
    main()